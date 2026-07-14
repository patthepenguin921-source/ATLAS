"""Document ingestion — turns uploaded files into searchable knowledge.

Pipeline: extract text → chunk → embed → store chunks (pgvector). Optionally
the Archivist agent enriches the document with a summary, keywords, and concept
links. Nothing is lost; everything becomes semantic memory.
"""
from __future__ import annotations

import io
from typing import Any

from app.core.supabase_client import eq, supabase
from app.embeddings.embedder import embed_texts

CHUNK_CHARS = 1400
CHUNK_OVERLAP = 200

IMAGE_EXTS = (".png", ".jpg", ".jpeg", ".heic", ".heif", ".webp", ".gif", ".bmp", ".tiff")


def is_image(mime_type: str, filename: str = "") -> bool:
    mt = (mime_type or "").lower()
    name = (filename or "").lower()
    return mt.startswith("image/") or name.endswith(IMAGE_EXTS)


def convert_image_to_pdf(content: bytes, filename: str = "") -> tuple[bytes, str]:
    """Convert an uploaded photo (png/jpg/heic/…) into a single-page PDF.

    Students photograph handouts and worksheets; we normalize every image into
    a PDF so it lives alongside the rest of their documents. HEIC (iPhone
    photos) needs pillow-heif's opener registered before Pillow can read it.
    Returns (pdf_bytes, pdf_filename).
    """
    import io

    from PIL import Image

    try:  # register HEIF/HEIC support if available
        from pillow_heif import register_heif_opener

        register_heif_opener()
    except Exception:
        pass

    img = Image.open(io.BytesIO(content))
    # img2pdf can't embed alpha/paletted images, so flatten onto white RGB.
    if img.mode in ("RGBA", "LA", "P"):
        background = Image.new("RGB", img.size, (255, 255, 255))
        img = img.convert("RGBA")
        background.paste(img, mask=img.split()[-1])
        img = background
    elif img.mode != "RGB":
        img = img.convert("RGB")

    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=90)

    import img2pdf

    pdf_bytes = img2pdf.convert(buf.getvalue())
    base = (filename or "photo").rsplit(".", 1)[0] or "photo"
    return pdf_bytes, f"{base}.pdf"


def ocr_image(content: bytes) -> str:
    """Best-effort OCR of a photo/scan so its text becomes searchable.

    Uses Tesseract via pytesseract. Degrades gracefully to an empty string if
    the OCR engine isn't installed or the image can't be read, so uploads never
    fail on account of OCR.
    """
    import io

    try:
        import pytesseract
        from PIL import Image

        try:  # HEIC support for iPhone photos
            from pillow_heif import register_heif_opener

            register_heif_opener()
        except Exception:
            pass

        img = Image.open(io.BytesIO(content))
        return (pytesseract.image_to_string(img) or "").strip()
    except Exception:
        return ""


# --------------------------------------------------------------------------
# Text extraction
# --------------------------------------------------------------------------
def _sanitize_text(text: str) -> str:
    """Strip NUL bytes and other chars Postgres `text` columns reject."""
    return text.replace("\x00", "")


def extract_text(content: bytes, mime_type: str, filename: str = "") -> str:
    name = (filename or "").lower()
    mt = (mime_type or "").lower()
    try:
        if "pdf" in mt or name.endswith(".pdf"):
            return _sanitize_text(_extract_pdf(content))
        if "presentation" in mt or name.endswith((".pptx", ".ppt")):
            return _sanitize_text(_extract_pptx(content))
    except Exception:
        pass  # fall through to plain-text decode
    try:
        return _sanitize_text(content.decode("utf-8", errors="ignore"))
    except Exception:
        return ""


def _extract_pdf(content: bytes) -> str:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(content))
    return "\n\n".join((page.extract_text() or "") for page in reader.pages)


def _extract_pptx(content: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(content))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts = [sh.text for sh in slide.shapes if getattr(sh, "has_text_frame", False)]
        if texts:
            parts.append(f"[Slide {i}] " + "\n".join(texts))
    return "\n\n".join(parts)


# --------------------------------------------------------------------------
# Chunking
# --------------------------------------------------------------------------
def chunk_text(text: str, size: int = CHUNK_CHARS, overlap: int = CHUNK_OVERLAP) -> list[str]:
    text = text.strip()
    if not text:
        return []
    chunks, start = [], 0
    while start < len(text):
        end = min(start + size, len(text))
        # try to break on a paragraph/sentence boundary
        if end < len(text):
            for sep in ("\n\n", "\n", ". "):
                cut = text.rfind(sep, start + size // 2, end)
                if cut != -1:
                    end = cut + len(sep)
                    break
        chunks.append(text[start:end].strip())
        if end >= len(text):
            break
        start = max(start + 1, end - overlap)
    return [c for c in chunks if c]


# --------------------------------------------------------------------------
# Ingest
# --------------------------------------------------------------------------
async def ingest_document(document_id: str, user_id: str, text: str) -> dict[str, Any]:
    """Chunk + embed + persist. Returns a small ingestion report."""
    text = _sanitize_text(text)

    # clear any prior chunks (re-ingest safe)
    await supabase.delete("document_chunks", filters={"document_id": eq(document_id)})

    chunks = chunk_text(text)
    if not chunks:
        await supabase.update(
            "documents",
            {"ingested": True, "ingest_error": "no extractable text", "extracted_text": text[:200000]},
            filters={"id": eq(document_id)},
        )
        return {"chunks": 0}

    embeddings = await embed_texts(chunks)
    rows = [
        {
            "user_id": user_id,
            "document_id": document_id,
            "chunk_index": i,
            "content": c,
            "token_count": max(1, len(c) // 4),
            "embedding": emb,
        }
        for i, (c, emb) in enumerate(zip(chunks, embeddings))
    ]
    # insert in batches to keep payloads reasonable
    for i in range(0, len(rows), 50):
        await supabase.insert("document_chunks", rows[i : i + 50])

    await supabase.update(
        "documents",
        {"ingested": True, "ingest_error": None, "extracted_text": text[:200000],
         "page_count": None},
        filters={"id": eq(document_id)},
    )
    return {"chunks": len(rows)}
